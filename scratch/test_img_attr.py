from pptx import Presentation
prs = Presentation(r"c:\Github repos\projects\Test Nexus\Test Nexus Proposition.pptx")
slide = prs.slides[1]
for shape in slide.shapes:
    if shape.shape_type == 13:
        rId = shape._element.blipFill.blip.rEmbed
        rel = slide.part.rels[rId]
        print("Rel type:", type(rel))
        print("Target part type:", type(rel.target_part))
        print("Target part dir:", dir(rel.target_part))
        break
