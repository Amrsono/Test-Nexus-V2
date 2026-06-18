import sys
from pptx import Presentation

def move_slide(prs, old_index, new_index):
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    slide_to_move = slides[old_index]
    sldIdLst.remove(slide_to_move)
    sldIdLst.insert(new_index, slide_to_move)

def add_bullet_points(tf, points):
    # tf is the text frame
    # Remove default empty paragraph if exists
    if len(tf.paragraphs) > 0 and tf.paragraphs[0].text == "":
        p0 = tf.paragraphs[0]
    else:
        p0 = tf.add_paragraph()
        
    for idx, (title, desc) in enumerate(points):
        if idx == 0 and p0.text == "":
            p = p0
        else:
            p = tf.add_paragraph()
        p.space_after = 12
        p.level = 0
        
        # Add bold title
        run_title = p.add_run()
        run_title.text = title + ": "
        run_title.font.bold = True
        
        # Add normal description
        run_desc = p.add_run()
        run_desc.text = desc

def main():
    pptx_path = r"c:\Github repos\projects\Test Nexus\Test Nexus Proposition.pptx"
    prs = Presentation(pptx_path)
    
    # We want to insert 3 slides between Slide 2 (index 1) and Slide 3 (index 2).
    # Layout 3: Two Content
    layout_two_content = prs.slide_layouts[3]
    
    # ---------------- SLIDE 3: The Challenge ----------------
    s3 = prs.slides.add_slide(layout_two_content)
    s3.shapes.title.text = "The Challenge: Testing Bottlenecks & Overhead"
    
    # Left column: Inefficient Workflows
    left_ph_s3 = s3.placeholders[1]
    left_tf_s3 = left_ph_s3.text_frame
    left_tf_s3.word_wrap = True
    add_bullet_points(left_tf_s3, [
        ("Legacy Spreadsheets & Docs", "Fragmented test plans and unstructured specifications stored in silos."),
        ("Manual Data Entry & Mapping", "Testers spend hours mapping Excel cells, copying test steps, and configuring tools manually."),
        ("Slow Execution Feedback", "Lack of instant alignment leads to late-stage bug detection and high risk of missed coverage.")
    ])
    
    # Right column: Capacity & Visibility
    right_ph_s3 = s3.placeholders[2]
    right_tf_s3 = right_ph_s3.text_frame
    right_tf_s3.word_wrap = True
    add_bullet_points(right_tf_s3, [
        ("Capacity Blindspots", "Difficulty tracking real-time tester capacity, leading to work imbalances and team burnout."),
        ("Reactive Risk Management", "No active warning signs for project slippage, resulting in delayed release timelines."),
        ("Disconnected Blocker Tracking", "Blockers and critical defects aren't dynamically linked to affected test coverage.")
    ])
    
    # ---------------- SLIDE 4: Core Capabilities - Parser & Capacity ----------------
    s4 = prs.slides.add_slide(layout_two_content)
    s4.shapes.title.text = "Core Capabilities: Speed & Smart Orchestration"
    
    # Left column: AI Parser & Mapper
    left_ph_s4 = s4.placeholders[1]
    left_tf_s4 = left_ph_s4.text_frame
    left_tf_s4.word_wrap = True
    add_bullet_points(left_tf_s4, [
        ("Zero-Mapping AI Parser", "Upload unstructured Excel test plans or word docs directly without pre-formatting."),
        ("Intelligent Field Mapping", "AI automatically detects, extracts, and maps IDs, summaries, steps, and expected results."),
        ("Instant Case Generation", "Converts plain text descriptions and messy documents into structured, clean test cases in seconds."),
        ("Apply-to-All Scenario Editor", "Bulk-edit scenario details and propagate changes instantly across all test steps with one click.")
    ])
    
    # Right column: Capacity Planner
    right_ph_s4 = s4.placeholders[2]
    right_tf_s4 = right_ph_s4.text_frame
    right_tf_s4.word_wrap = True
    add_bullet_points(right_tf_s4, [
        ("Automated Capacity Planning", "Calculates daily workloads automatically based on remaining test cases and target sprint dates."),
        ("Dynamic Auto-Assignment", "AI balances the load across the team, auto-assigning cases to available testers."),
        ("Adjustable Validation Points", "Define custom verification steps and edit validation criteria dynamically to fit changing requirements."),
        ("Manual Control & Override", "Easily adjust assignments and re-assign blocked cases through a drag-and-drop panel.")
    ])
    
    # ---------------- SLIDE 5: Core Capabilities - AI Advisor & Metrics ----------------
    s5 = prs.slides.add_slide(layout_two_content)
    s5.shapes.title.text = "Project Intelligence: AI Advisor & Live Dashboards"
    
    # Left column: AI Advisor Agent
    left_ph_s5 = s5.placeholders[1]
    left_tf_s5 = left_ph_s5.text_frame
    left_tf_s5.word_wrap = True
    add_bullet_points(left_tf_s5, [
        ("Proactive Risk Detection", "Alerts if progress velocity drops below 90% or if failure rates exceed 20% in specific modules."),
        ("Daily Actionable Advice", "Generates daily strategic guidance (e.g., 'Prioritize blocker #402 to unlock 15 dependent test cases')."),
        ("Configurable Gemini API", "Powered by Google Gemini through API connections configurable in app settings to swap or adjust models."),
        ("Reports & Burndown Charts", "Comprehensive analytics showing sprint-level burndowns, cycle times, and bug-resolution velocity.")
    ])
    
    # Right column: Multi-Project Workspace
    right_ph_s5 = s5.placeholders[2]
    right_tf_s5 = right_ph_s5.text_frame
    right_tf_s5.word_wrap = True
    add_bullet_points(right_tf_s5, [
        ("Tabbed Multi-Project Dashboard", "Monitor and toggle between multiple active projects seamlessly using navigation tabs."),
        ("Custom Identity & Themes", "Personalize each project tab with its own logo, custom name, and distinct theme colors."),
        ("Multi-Tenant Architecture", "Securely partition projects, logs, and configurations across multiple tenants and client workspaces."),
        ("Live Execution Charts", "Real-time burndown charts and blocker hubs sync instantly as testers log pass/fail status.")
    ])

    # Now we have slides:
    # 0: Slide 1 (Title)
    # 1: Slide 2 (Architecture Diagram)
    # 2: Slide 3 (Demo Slide, old index 2)
    # 3: Slide 4 (Video Slide, old index 3)
    # 4: Slide 5 (Blank Slide, old index 4)
    # 5: Slide 6 (The Challenge, new)
    # 6: Slide 7 (Parser/Capacity, new)
    # 7: Slide 8 (Advisor/Metrics, new)
    
    # Let's delete the blank slide (old Slide 5, which is now index 4)
    id_list = prs.slides._sldIdLst
    slide_id = id_list[4]
    id_list.remove(slide_id)
    # Remove from relationship
    prs.part.drop_rel(slide_id.rId)
    
    # Now we have slides:
    # 0: Slide 1 (Title)
    # 1: Slide 2 (Architecture Diagram)
    # 2: Slide 3 (Demo Slide)
    # 3: Slide 4 (Video Slide)
    # 4: Slide 5 (The Challenge)
    # 5: Slide 6 (Parser/Capacity)
    # 6: Slide 7 (Advisor/Metrics)
    
    # Let's move the new slides:
    # The Challenge (index 4) -> move to index 2
    # Parser/Capacity (index 5) -> move to index 3
    # Advisor/Metrics (index 6) -> move to index 4
    move_slide(prs, 4, 2)
    move_slide(prs, 5, 3)
    move_slide(prs, 6, 4)
    
    prs.save(pptx_path)
    print("PowerPoint deck populated successfully.")

if __name__ == "__main__":
    main()
