from pptx import Presentation

prs = Presentation(r"c:\Github repos\projects\Test Nexus\Test Nexus Proposition.pptx")
for i, slide in enumerate(prs.slides):
    print(f"\n================ SLIDE {i+1} ================")
    for j, shape in enumerate(slide.shapes):
        print(f"Shape {j+1}: Name='{shape.name}', ID={shape.shape_id}, Type={shape.shape_type}")
        if shape.has_text_frame:
            for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                p_text = "".join(run.text for run in paragraph.runs)
                if p_text.strip():
                    print(f"  P{p_idx+1}: {p_text}")
        if shape.has_table:
            table = shape.table
            print("  Table:")
            for r_idx, row in enumerate(table.rows):
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text.strip())
                print(f"    Row {r_idx+1}: {row_text}")
