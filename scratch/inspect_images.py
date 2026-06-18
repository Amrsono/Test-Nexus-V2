from pptx import Presentation
import os

prs = Presentation(r"c:\Github repos\projects\Test Nexus\Test Nexus Proposition.pptx")
print("Images in slides:")
for i, slide in enumerate(prs.slides):
    for j, shape in enumerate(slide.shapes):
        if shape.shape_type == 13: # PICTURE
            image = shape.image
            print(f"Slide {i+1}, Shape {j+1}: Name='{shape.name}', Size={len(image.blob)} bytes, Content Type={image.content_type}")
            # Save the image to inspect it or locate where it is
            filename = f"extracted_image_s{i+1}_sh{j+1}.png"
            with open(filename, 'wb') as f:
                f.write(image.blob)
            print(f"  Saved to {filename}")
