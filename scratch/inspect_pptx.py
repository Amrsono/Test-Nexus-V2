import sys
from pptx import Presentation

def inspect_presentation(file_path):
    prs = Presentation(file_path)
    print(f"Total slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        if slide.shapes.title:
            print(f"Title: {slide.shapes.title.text}")
        else:
            print("No Title Shape")
        
        # Check all text frames
        for j, shape in enumerate(slide.shapes):
            shape_type = shape.shape_type
            name = shape.name
            print(f"  Shape {j+1}: Name='{name}', Type={shape_type}")
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    print(f"    Text: {text[:200]}...")
            if shape.has_table:
                print("    Contains Table")
                table = shape.table
                for r_idx, row in enumerate(table.rows):
                    row_text = [cell.text.strip() for cell in row.cells]
                    print(f"      Row {r_idx+1}: {row_text}")

if __name__ == "__main__":
    file_path = r"c:\Github repos\projects\Test Nexus\Test Nexus Proposition.pptx"
    inspect_presentation(file_path)
