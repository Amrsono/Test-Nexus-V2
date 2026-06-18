from pptx import Presentation

def main():
    pptx_path = r"c:\Github repos\projects\Test Nexus\Test Nexus Proposition.pptx"
    new_image_path = r"C:\Users\Sono\.gemini\antigravity\brain\8628bbd8-3570-4e00-9d74-c52abeec1dba\gemini_test_nexus_logo_1781541436987.png"
    
    prs = Presentation(pptx_path)
    slide2 = prs.slides[1]
    
    picture_shape = None
    for shape in slide2.shapes:
        if shape.shape_type == 13: # PICTURE
            picture_shape = shape
            break
            
    if picture_shape:
        rId = picture_shape._element.blipFill.blip.rEmbed
        rel = slide2.part.rels[rId]
        print(f"Replacing image in target part: {rel.target_part.partname}")
        with open(new_image_path, 'rb') as f:
            new_blob = f.read()
        
        # Overwrite the part's internal blob
        rel.target_part._blob = new_blob
        
        prs.save(pptx_path)
        print("Image replaced successfully.")
    else:
        print("Picture shape not found.")

if __name__ == "__main__":
    main()
